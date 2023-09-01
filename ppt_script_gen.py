import os
import time
import traceback
import json
from pptx import Presentation


def load_json(fdirname):
    with open(fdirname, 'r') as file:
        btc = json.load(file)
    return btc
def save_json(dic, fdirname):
    with open(fdirname, 'w') as file:
        json.dump(dic, file, indent=4)
def save_txt(txt, fdirname):
    with open(fdirname, 'w') as file:
        file.write(txt)
def load_txt(fdirname):
    with open(fdirname, 'r') as file:
        txt = file.read().strip()
    return txt
 
def summary_page_content(shapes, page):
    def normalize(value, max_value):
        return value / max_value

    sumtext = "Page: {}\n".format(page)
    
    if shapes['texts']:
        max_width = max([shape['width'] for shape in shapes['texts']])
        max_height = max([shape['height'] for shape in shapes['texts']])
    else:
        max_width = max_height = 1  # Default values to avoid division by zero
    
    for shape_type, shape_info_list in shapes.items():
        if shape_info_list:
            sumtext += "{}:\n".format(shape_type)
            for info in shape_info_list:
                norm_left = (normalize(info['left'], max_width))
                norm_top = (normalize(info['top'], max_height))
                norm_width = (normalize(info['width'], max_width))
                norm_height = (normalize(info['height'], max_height))
                
                if shape_type == 'texts' and info['text']:
                    sumtext += "  Position: ({:.3f}, {:.3f})\n".format(norm_left, norm_top)
                    sumtext += "  Size: ({:.3f} x {:.3f})\n".format(norm_width, norm_height)
                    sumtext += "  Text: {}\n".format(info['text'])
                elif shape_type == 'images':
                    sumtext += "  Position: ({:.3f}, {:.3f})\n".format(norm_left, norm_top)
                    sumtext += "  Size: ({:.3f} x {:.3f})\n".format(norm_width, norm_height)
                    if 'image_info' in info:
                        sumtext += "  Image: {}\n".format(info['image_info'])
                elif shape_type == 'tables':
                    sumtext += "  Position: ({}, {})\n".format(norm_left, norm_top)
                    sumtext += "  Size: ({:.3f} x {:.3f})\n".format(norm_width, norm_height)
                    sumtext += "  Table:\n"
                    if 'rows' in info:
                        for row in info['rows']:
                            sumtext += "    {}\n".format(row)
                    else:
                        sumtext += "    No table data available\n"
                elif shape_type == 'charts':
                    sumtext += "  Position: ({}, {})\n".format(norm_left, norm_top)
                    sumtext += "  Size: ({:.3f} x {:.3f})\n".format(norm_width, norm_height)
                    if 'chart_info' in info:
                        sumtext += "  Chart: {}\n".format(info['chart_info'])
                    else:
                        sumtext += "    No chart data available\n"
                sumtext += "\n"  # Add an extra newline after each instance
    return sumtext

def read_ppt(slide):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    res = {'texts': [], 'images': [], 'tables': [], 'charts': []}
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            text = "\n".join([p.text for p in text_frame.paragraphs])
            
            shape_info = {
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'text': text
            }
            res['texts'].append(shape_info)
        
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_info = {
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'image_path': shape.image.blob
            }
            res['images'].append(image_info)
            
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_info = {
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'rows': []
            }
            
            for row in shape.table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = "\n".join([p.text for p in cell.text_frame.paragraphs])
                    row_text.append(cell_text)
                table_info['rows'].append(row_text)
            
            res['tables'].append(table_info)
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            chart_info = {
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'chart_type': shape.chart.chart_type,
                'chart_title': shape.chart.chart_title.text_frame.text if shape.chart.chart_title is not None else ""
            }
            res['charts'].append(chart_info)

    return res

def callGPT_API(messages):
    import openai
    openai.api_key = 'your-api-key'
    itry = 0
    while itry<3:
        try:
            response = openai.ChatCompletion.create(model="gpt-3.5-turbo", messages=[{'role': 'user', 'content': messages}])

            return response.choices[0].message.content.strip()
        except:
            print(traceback.format_exc())
            time.sleep(1)
            itry += 1
            print('error occered in call gpt, tried {} times'.format(itry))
            pass
    return 'errored too many times???'


def auto_summary_ppt(abs_file, save_path, sentence_cnt, use_paid_API=False):
    content_str_lst = load_json(f'{save_path}/layouts.json')
    speech_texts = ''
    new_gen = ''

    for page, content_str in enumerate(content_str_lst):
        abss = load_txt(abs_file)
        
        prev_pages = max(page - 2, 0)
        next_pages = min(page + 3, len(content_str_lst))
        context_content = "\n".join(content_str_lst[prev_pages:next_pages])  # Context includes 5 pages (prev 2, current, next 2)

        if new_gen:
            prev = 'page {} / {}:\n{}'.format(page, len(content_str_lst), new_gen) 
        else:
            prev = 'Not applicable. Current slide is the beginning page.'

        prompt = '''
Please write a script of speech presentation, based on the powerpoint slide layouts. I'll provide you with the presentation [abstract], the [layout] of [current] slide, previous slides, and later slides, and the [previous slide speech script], if any. 
Please generate content only for the [current] slide, while considering the context of the previous and later slides to make it coherent. Unless it is the first slides, do NOT begin with words like 'Ladies and gentlemen' -- no one say this in the middle of presentation.

[Abstract]: 
{}
The page [layouts]: 
{}
The [current] slide page is: 
{}
[previous slide speech script]:
{}


Please limit to less than or equal to {} sentences. Please limit your word/sentence count in the this way: 
- For each sentence in the script, use a new line.
- For each new line, begin with an incrementally increasing line number, e.g, 1,2, ..., start with 1.
- Control your presentation progress by looking at current sentence count.
- You should finish before you reach the sentence count upper limit of {}.

Please generate the current page speech script now. Please directly write results, do not analyze, and do not say any confirmation words to me like 'OK I understand', etc.
'''.format(abss, 
           context_content, 
           '{} / {}'.format(page+1, len(content_str_lst)), 
           prev,
           sentence_cnt,
           sentence_cnt,
           )
        new_gen = callGPT_API(prompt) if use_paid_API else ''
        save_txt(prompt, f'{save_path}/prompt-{page}.txt')
        speech_texts = speech_texts + '\n-------------\nPage {} / {}:\n'.format(page+1, len(content_str_lst)) + new_gen
    
    save_fdn = f'{save_path}/chatGPT_API_result.txt'
    save_txt(speech_texts, save_fdn)
    return speech_texts



def summarize_layout(pptx_path, save_path):
    presentation = Presentation(pptx_path)

    total_pages = len(presentation.slides)
    one_ppt = []
    for page, slide in enumerate(presentation.slides):
        shapes = read_ppt(slide)
        sumtext = summary_page_content(shapes, '{} / {}'.format(page + 1, total_pages))
        one_ppt.append(sumtext)
        print(sumtext, file=open(f'{save_path}/layouts-{page}.txt', 'w'))
    save_json(one_ppt, f'{save_path}/layouts.json')


def main():
    pptx_path = 'demo.pptx'
    abs_file = 'abs.txt'
    save_path = 'results'
    sentence_cnt = 6  # how many sentences you want chatGPT to speak for one page.
    os.makedirs(save_path, exist_ok=1)
    summarize_layout(pptx_path, save_path)
    auto_summary_ppt(abs_file, save_path, sentence_cnt, use_paid_API=False)
    return


if __name__ == '__main__':
    main()